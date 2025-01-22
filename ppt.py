from enum import Enum, IntEnum
from typing import Annotated
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml, CT_OfficeStyleSheet, CT_SRgbColor
from lxml import etree
from pydantic import BaseModel, Field, StringConstraints


HexColor = Annotated[
    str, StringConstraints(pattern=r"^#([A-F0-9]{6})$", min_length=7, max_length=7)
]


class ColorChoice(BaseModel):
    color: HexColor = Field(..., description="A 6-digit Hex code in uppercase")
    reason: str


class FontChoice(BaseModel):
    family: str = Field(..., description="The name of the font")
    reason: str


class ThemeFonts(BaseModel):
    header: FontChoice
    body: FontChoice


class ThemeColors(BaseModel):
    dark: ColorChoice
    light: ColorChoice
    accent1: ColorChoice
    accent2: ColorChoice
    accent3: ColorChoice
    accent4: ColorChoice
    accent5: ColorChoice
    accent6: ColorChoice
    hyperlink: ColorChoice
    followed_hyperlink: ColorChoice


class Theme(BaseModel):
    colors: ThemeColors
    fonts: ThemeFonts


class Schema(BaseModel):
    background: ColorChoice
    theme: Theme


TRANSFORMATIVE_COLOR_MAP = {
    "dark": "dk2",
    "light": "lt2",
    "hyperlink": "hlink",
    "followed_hyperlink": "folHlink",
}


def theme_colors_to_pptx_format(theme_colors: ThemeColors) -> dict[str, str]:
    return {
        TRANSFORMATIVE_COLOR_MAP.get(name, name): color_choice['color'].lstrip("#")
        for name, color_choice in theme_colors.model_dump().items()
    }


namespaces = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def set_fonts(theme_xml: CT_OfficeStyleSheet, header_font: str, body_font: str):
    font_root = theme_xml.find(".//a:fontScheme", namespaces=namespaces)
    if not isinstance(font_root, etree._Element):
        raise ValueError("Font scheme is not a _Element")

    def set_font(font_xpath: str, font_name: str):
        font_elem = font_root.find(font_xpath, namespaces=namespaces)
        if not isinstance(font_elem, etree._Element):
            raise ValueError(f"Font element at {font_xpath} is not a _Element")

        latin = font_elem.find(".//a:latin", namespaces=namespaces)
        if not isinstance(latin, etree._Element):
            raise ValueError("latin is not a _Element")
        latin.set("typeface", font_name)

    set_font(".//a:majorFont", header_font)
    set_font(".//a:minorFont", body_font)


def get_color_scheme_elem(theme_xml: CT_OfficeStyleSheet, key: str):
    elem = theme_xml.find(
        f".//a:themeElements/a:clrScheme/a:{key}", namespaces=namespaces
    )
    if not isinstance(elem, etree._Element):
        raise ValueError(f"{key} is not a _Element")
    return elem


def modify_secondary_color(theme_xml: CT_OfficeStyleSheet, key: str, color: str):
    elem = get_color_scheme_elem(theme_xml, key)
    srgb_clr = elem.find(".//a:srgbClr", namespaces=namespaces)
    if not isinstance(srgb_clr, CT_SRgbColor):
        raise ValueError("srgbClr is not a _Element")
    srgb_clr.set("val", color)


class SlideLayout(IntEnum):
    TITLE = 0
    TITLE_AND_CONTENT = 1


def create_ppt(response_schema: Schema, output_path: str):
    prs = Presentation()
    slide_master = prs.slide_master
    background_fill = slide_master.background.fill
    background_fill.solid()
    background_fill.fore_color.rgb = RGBColor.from_string(
        response_schema.background.color.lstrip("#")
    )

    theme_part = slide_master.part.part_related_by(RELATIONSHIP_TYPE.THEME)
    theme_xml = parse_xml(theme_part.blob)
    if not isinstance(theme_xml, CT_OfficeStyleSheet):
        raise ValueError("Theme part is not a CT_OfficeStyleSheet")

    test_dict = theme_colors_to_pptx_format(response_schema.theme.colors)
    for key, color in test_dict.items():
        try:
            modify_secondary_color(theme_xml, key, color)
        except ValueError as e:
            print(e)

    set_fonts(
        theme_xml,
        response_schema.theme.fonts.header.family,
        response_schema.theme.fonts.body.family,
    )
    modified_theme_blobs = etree.tostring(
        theme_xml, xml_declaration=True, encoding="UTF-8"
    )
    theme_part._blob = modified_theme_blobs

    slide_layout = slide_master.slide_layouts[SlideLayout.TITLE]
    title_slide = prs.slides.add_slide(slide_layout)
    title_slide.shapes.title.text = "PLACEHOLDER TITLE"

    prs.save(output_path)


if __name__ == "__main__":
    schema = Schema.model_validate(
        {
            "background": {"color": "#FFFFFF", "reason": "White background"},
            "theme": {
                "colors": {
                    "dark": {"color": "#000000", "reason": "Black text"},
                    "light": {"color": "#FFFFFF", "reason": "White background"},
                    "accent1": {"color": "#0000FF", "reason": "Blue accent"},
                    "accent2": {"color": "#FF0000", "reason": "Red accent"},
                    "accent3": {"color": "#00FF00", "reason": "Green accent"},
                    "accent4": {"color": "#FFFF00", "reason": "Yellow accent"},
                    "accent5": {"color": "#00FFFF", "reason": "Cyan accent"},
                    "accent6": {"color": "#FF00FF", "reason": "Magenta accent"},
                    "hyperlink": {"color": "#0000FF", "reason": "Blue hyperlink"},
                    "followed_hyperlink": {
                        "color": "#800080",
                        "reason": "Purple followed hyperlink",
                    },
                },
                "fonts": {
                    "header": {"typeface": "Comic Sans MS", "reason": "Arial for headers"},
                    "body": {"typeface": "Comic Sans MS", "reason": "Arial for body text"},
                },
            },
        }
    )

    create_ppt(schema, "C:/Users/Glizzus/Documents/dev_container_template.pptx")
